# -*- coding: utf-8 -*-
"""
迪士尼RAG助手 - 知识库索引构建

功能：递归解析知识库文档/图片/视频，生成embedding，构建FAISS索引并保存
"""
import os
import base64
import json
import subprocess
import numpy as np
import faiss
import dashscope
from http import HTTPStatus
from docx import Document as DocxDocument
from win32com_doc_ppt_parser import parse_doc_with_win32com, parse_ppt_with_win32com
from pptx import Presentation
from pptx.util import Inches
import pdfplumber

# 配置
DASHSCOPE_API_KEY = os.getenv("DASHSCOPE_API_KEY2")
if not DASHSCOPE_API_KEY:
    raise ValueError("错误：请设置 'DASHSCOPE_API_KEY2' 环境变量。")

dashscope.api_key = DASHSCOPE_API_KEY

DOCS_DIR = "迪士尼RAG知识库（完整）"
MULTIMODAL_EMBEDDING_MODEL = "tongyi-embedding-vision-plus"
CHUNK_SIZE = 500
CHUNK_OVERLAP = 50
INDEX_FILE = "disney_index.faiss"
METADATA_FILE = "disney_metadata.json"

# 视频知识库
VIDEO_KNOWLEDGE = [
    {
        "url": "https://dataset-1255932437.cos.ap-nanjing.myqcloud.com/mp4/car.mp4",
        "description": "汽车剐蹭视频"
    }
]

IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp'}


def is_ole2_format(file_path):
    """读取文件前8字节判断是否为OLE2旧版二进制格式"""
    with open(file_path, 'rb') as f:
        header = f.read(8)
    return header[:4] == b'\xd0\xcf\x11\xe0'


def is_ooxml_format(file_path):
    """判断是否为OOXML格式（PK开头）"""
    with open(file_path, 'rb') as f:
        header = f.read(2)
    return header[:2] == b'PK'


def collect_all_files(root_dir):
    """递归收集目录下所有文件，返回(相对路径, 绝对路径)列表"""
    files = []
    for dirpath, dirnames, filenames in os.walk(root_dir):
        for filename in filenames:
            if filename.startswith('.') or filename.startswith('~$'):
                continue
            abs_path = os.path.join(dirpath, filename)
            rel_path = os.path.relpath(abs_path, root_dir)
            files.append((rel_path, abs_path))
    return files


def parse_docx(file_path):
    """解析DOCX文件，提取段落和表格（Markdown格式）"""
    doc = DocxDocument(file_path)
    all_text = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            all_text.append(text)

    for table in doc.tables:
        if table.rows:
            md_table = []
            header = [cell.text.strip() for cell in table.rows[0].cells]
            md_table.append("| " + " | ".join(header) + " |")
            md_table.append("|" + "|".join(["---"] * len(header)) + "|")
            for row in table.rows[1:]:
                row_data = [cell.text.strip() for cell in row.cells]
                md_table.append("| " + " | ".join(row_data) + " |")
            all_text.append("\n".join(md_table))

    return "\n".join(all_text)


def parse_doc_with_antiword(file_path):
    """使用antiword命令行工具提取旧版.doc文本"""
    try:
        result = subprocess.run(
            ['antiword', '-m', 'UTF-8', file_path],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode != 0:
            print(f"    警告: antiword 失败: {result.stderr.strip()}")
            return None
        if not result.stdout:
            print(f"    警告: antiword 返回空内容，跳过")
            return None
        text = result.stdout.strip()
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        return '\n'.join(lines)
    except FileNotFoundError:
        print(f"    警告: antiword 未安装，跳过旧版.doc文件")
        return None
    except subprocess.TimeoutExpired:
        print(f"    警告: antiword 超时，跳过该文件")
        return None


def parse_pptx(file_path):
    """解析PPTX文件，提取幻灯片文本、表格和备注"""
    prs = Presentation(file_path)
    all_text = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_content.append(text)
            if shape.has_table:
                table = shape.table
                rows_list = list(table.rows)
                if rows_list:
                    md_table = []
                    header = [cell.text.strip() for cell in rows_list[0].cells]
                    md_table.append("| " + " | ".join(header) + " |")
                    md_table.append("|" + "|".join(["---"] * len(header)) + "|")
                    for row in rows_list[1:]:
                        row_data = [cell.text.strip() for cell in row.cells]
                        md_table.append("| " + " | ".join(row_data) + " |")
                    slide_content.append("\n".join(md_table))

        slide_text = f"[幻灯片 {slide_idx}] " + "\n".join(slide_content)

        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                notes_text = f"\n[备注] {notes}"

        all_text.append(slide_text + notes_text)

    return "\n".join(all_text)


def parse_pdf(file_path):
    """使用pdfplumber逐页提取PDF文本"""
    all_text = []
    with pdfplumber.open(file_path) as pdf:
        for page_idx, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                all_text.append(f"[第{page_idx}页] {text.strip()}")
    return "\n".join(all_text)


def split_text(text, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    """按固定窗口滑动切分文本"""
    if not text or not text.strip():
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap
    return chunks


def get_text_embedding(text):
    """文本embedding"""
    resp = dashscope.MultiModalEmbedding.call(
        model=MULTIMODAL_EMBEDDING_MODEL,
        input=[{'text': text}]
    )
    if resp.status_code != HTTPStatus.OK:
        raise Exception(f"文本Embedding失败: {resp.message}")
    return resp.output['embeddings'][0]['embedding']


def get_image_embedding(image_path):
    """图片embedding"""
    with open(image_path, "rb") as f:
        base64_image = base64.b64encode(f.read()).decode('utf-8')

    ext = os.path.splitext(image_path)[1].lower().lstrip('.')
    if ext == 'jpg':
        ext = 'jpeg'
    image_data = f"data:image/{ext};base64,{base64_image}"

    resp = dashscope.MultiModalEmbedding.call(
        model=MULTIMODAL_EMBEDDING_MODEL,
        input=[{'image': image_data}]
    )
    if resp.status_code != HTTPStatus.OK:
        raise Exception(f"图片Embedding失败: {resp.message}")
    return resp.output['embeddings'][0]['embedding']


def get_video_embedding(video_url):
    """视频embedding（多帧取平均）"""
    resp = dashscope.MultiModalEmbedding.call(
        model=MULTIMODAL_EMBEDDING_MODEL,
        input=[{'video': video_url}]
    )
    if resp.status_code != HTTPStatus.OK:
        raise Exception(f"视频Embedding失败: {resp.message}")

    embeddings = resp.output['embeddings']
    if len(embeddings) > 1:
        vectors = [np.array(e['embedding']) for e in embeddings]
        return np.mean(vectors, axis=0).tolist()
    return embeddings[0]['embedding']


def build_and_save():
    """构建知识库并保存FAISS索引"""
    print("\n--- 构建多模态知识库 ---")
    print(f"知识库目录: {DOCS_DIR}")
    print(f"切分参数: chunk_size={CHUNK_SIZE}, overlap={CHUNK_OVERLAP}")

    if not os.path.isdir(DOCS_DIR):
        raise ValueError(f"错误：知识库目录不存在: {DOCS_DIR}")

    all_files = collect_all_files(DOCS_DIR)
    if not all_files:
        print("知识库目录为空，不生成索引文件。")
        return

    metadata_store = []
    all_vectors = []
    doc_id = 0

    for rel_path, abs_path in all_files:
        ext = os.path.splitext(rel_path)[1].lower()

        try:
            if ext == '.docx':
                print(f"  处理: {rel_path}")
                full_text = parse_docx(abs_path)
                if not full_text:
                    continue
                chunks = split_text(full_text)
                print(f"    {len(full_text)} 字符 → {len(chunks)} 个chunk")
                for chunk in chunks:
                    vector = get_text_embedding(chunk)
                    all_vectors.append(vector)
                    metadata_store.append({
                        "id": doc_id,
                        "source": rel_path,
                        "type": "text",
                        "content": chunk
                    })
                    doc_id += 1

            elif ext == '.doc':
                if is_ooxml_format(abs_path):
                    print(f"  处理 (OOXML→docx): {rel_path}")
                    full_text = parse_docx(abs_path)
                else:
                    print(f"  处理 (win32com): {rel_path}")
                    full_text = parse_doc_with_win32com(abs_path)
                    if full_text is None:
                        print(f"    降级: win32com 失败，尝试 antiword: {rel_path}")
                        full_text = parse_doc_with_antiword(abs_path)
                        if full_text is None:
                            continue
                if not full_text:
                    continue
                chunks = split_text(full_text)
                print(f"    {len(full_text)} 字符 → {len(chunks)} 个chunk")
                for chunk in chunks:
                    vector = get_text_embedding(chunk)
                    all_vectors.append(vector)
                    metadata_store.append({
                        "id": doc_id,
                        "source": rel_path,
                        "type": "text",
                        "content": chunk
                    })
                    doc_id += 1

            elif ext == '.pptx':
                print(f"  处理: {rel_path}")
                full_text = parse_pptx(abs_path)
                if not full_text:
                    continue
                chunks = split_text(full_text)
                print(f"    {len(full_text)} 字符 → {len(chunks)} 个chunk")
                for chunk in chunks:
                    vector = get_text_embedding(chunk)
                    all_vectors.append(vector)
                    metadata_store.append({
                        "id": doc_id,
                        "source": rel_path,
                        "type": "text",
                        "content": chunk
                    })
                    doc_id += 1

            elif ext == '.ppt':
                if is_ooxml_format(abs_path):
                    print(f"  处理 (OOXML→pptx): {rel_path}")
                    full_text = parse_pptx(abs_path)
                else:
                    print(f"  处理 (win32com): {rel_path}")
                    full_text = parse_ppt_with_win32com(abs_path)
                    if full_text is None:
                        print(f"    警告: win32com 解析失败，跳过: {rel_path}")
                        continue
                if not full_text:
                    continue
                chunks = split_text(full_text)
                print(f"    {len(full_text)} 字符 → {len(chunks)} 个chunk")
                for chunk in chunks:
                    vector = get_text_embedding(chunk)
                    all_vectors.append(vector)
                    metadata_store.append({
                        "id": doc_id,
                        "source": rel_path,
                        "type": "text",
                        "content": chunk
                    })
                    doc_id += 1

            elif ext == '.pdf':
                print(f"  处理: {rel_path}")
                full_text = parse_pdf(abs_path)
                if not full_text:
                    continue
                chunks = split_text(full_text)
                print(f"    {len(full_text)} 字符 → {len(chunks)} 个chunk")
                for chunk in chunks:
                    vector = get_text_embedding(chunk)
                    all_vectors.append(vector)
                    metadata_store.append({
                        "id": doc_id,
                        "source": rel_path,
                        "type": "text",
                        "content": chunk
                    })
                    doc_id += 1

            elif ext in IMAGE_EXTENSIONS:
                print(f"  处理图片: {rel_path}")
                vector = get_image_embedding(abs_path)
                all_vectors.append(vector)
                metadata_store.append({
                    "id": doc_id,
                    "source": rel_path,
                    "type": "image",
                    "path": rel_path,
                    "content": f"[图片] {os.path.basename(rel_path)}"
                })
                doc_id += 1

            else:
                print(f"  跳过未知格式: {rel_path}")

        except Exception as e:
            print(f"    错误: {rel_path} → {e}")
            continue

    # 处理视频知识库
    print("  处理视频知识库...")
    for video_info in VIDEO_KNOWLEDGE:
        print(f"    - {video_info['description']}")
        try:
            vector = get_video_embedding(video_info["url"])
            all_vectors.append(vector)
            metadata_store.append({
                "id": doc_id,
                "source": f"视频: {video_info['description']}",
                "type": "video",
                "url": video_info["url"],
                "description": video_info["description"],
                "content": f"[视频] {video_info['description']}"
            })
            doc_id += 1
        except Exception as e:
            print(f"    错误: {video_info['description']} → {e}")

    if not all_vectors:
        print("未生成任何向量，不生成索引文件。")
        return

    dim = len(all_vectors[0])
    print(f"\n向量维度: {dim}")
    print(f"总条目数: {len(all_vectors)}")

    index = faiss.IndexFlatL2(dim)
    index.add(np.array(all_vectors).astype('float32'))

    faiss.write_index(index, INDEX_FILE)
    print(f"索引已保存: {INDEX_FILE}")

    with open(METADATA_FILE, 'w', encoding='utf-8') as f:
        json.dump(metadata_store, f, ensure_ascii=False, indent=2)
    print(f"元数据已保存: {METADATA_FILE}")

    text_count = sum(1 for m in metadata_store if m["type"] == "text")
    image_count = sum(1 for m in metadata_store if m["type"] == "image")
    video_count = sum(1 for m in metadata_store if m["type"] == "video")
    print(f"\n完成! 文本:{text_count}, 图片:{image_count}, 视频:{video_count}")


if __name__ == "__main__":
    build_and_save()
